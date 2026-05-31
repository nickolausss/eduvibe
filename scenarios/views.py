import json
import logging
import os
import io
from datetime import datetime, date, timedelta
import re
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from django.http import JsonResponse, HttpResponse
from django.template.loader import render_to_string
from django.conf import settings
from django.db.models import Count
from .models import Scenario, ScenarioVersion, ScheduledTheme, LLMRequestLog
from .forms import ScenarioCreateForm
from .services import ScenarioGeneratorService, ScenarioVersionService

logger = logging.getLogger(__name__)

TYPE_TRANSLATIONS = {
    'org': 'Оргмомент',
    'challenge': 'Вызов',
    'main': 'Основная часть',
    'reflection': 'Рефлексия',
    'closing': 'Заключение',
}

KEY_TRANSLATIONS = {
    'educational': 'Воспитательная цель',
    'practical': 'Практическая цель',
    'soft_skills': 'Soft Skills',
    'meta_subject': 'Метапредметные результаты',
    'personal': 'Личностные результаты',
    'equipment': 'Оборудование',
    'stationery': 'Канцелярия',
    'digital': 'Цифровые ресурсы',
    'props': 'Реквизит',
    'handouts': 'Раздаточные материалы',
    'materials_full': 'Материалы и ресурсы',
    'for_juniors': 'Для младших классов',
    'for_seniors': 'Для старших классов',
    'for_ovz': 'Для детей с ОВЗ',
    'for_gifted': 'Для одарённых детей',
    'adaptation': 'Адаптация',
    'method': 'Метод проведения',
    'questions': 'Вопросы для рефлексии',
    'teacher_script': 'Слова учителя',
    'reflection': 'Рефлексия',
    'self_assessment': 'Самооценка',
    'group_reflection': 'Групповая рефлексия',
    'technical': 'Технические риски',
    'methodical': 'Методические риски',
    'dynamic': 'Риски динамики группы',
    'conflict': 'Конфликтные ситуации',
    'time_management': 'Тайминг',
    'low_engagement': 'Низкая вовлечённость',
    'quantitative': 'Количественная оценка',
    'qualitative': 'Качественная оценка',
    'criteria': 'Критерии',
    'indicators': 'Показатели',
    'teacher_tips': 'Советы учителю',
    'before_lesson': 'Перед занятием',
    'during_lesson': 'Во время занятия',
    'after_lesson': 'После занятия',
    'week_before': 'За неделю',
    'day_before': 'За день',
    'hour_before': 'За час',
    'five_minutes': 'За 5 минут',
    'checklist': 'Чек-лист подготовки',
    'homework': 'Домашнее задание',
    'variants': 'Варианты',
    'basic': 'Базовый уровень',
    'advanced': 'Продвинутый уровень',
    'creative': 'Творческое задание',
    'deadline': 'Срок выполнения',
    'name': 'Название',
    'duration_minutes': 'Длительность (мин)',
    'type': 'Тип этапа',
    'description': 'Описание',
    'teacher_actions': 'Действия учителя',
    'student_actions': 'Действия учеников',
    'materials': 'Материалы',
    'mechanics': 'Механика',
    'expected_result': 'Ожидаемый результат',
    'title': 'Название',
    'theme': 'Тема',
    'grade': 'Класс',
    'duration': 'Длительность',
    'format': 'Формат',
    'direction': 'Направление',
    'goals': 'Цели',
    'stages': 'Этапы',
    'legend': 'Легенда',
    'risks': 'Риски',
    'assessment': 'Оценка эффективности',
    'host_script': 'Сценарий для ведущего',
    'question': 'Вопрос',
    'text': 'Текст',
    'expected_answers': 'Ожидаемые ответы',
    'teacher_response': 'Реакция учителя',
    'answer': 'Ответ',
    'student': 'Ученик',
    'action': 'Действие',
    'quantity': 'Количество',
    'item': 'Предмет',
}


def translate_key(key):
    return KEY_TRANSLATIONS.get(key, key)


def get_client_ip(request):
    x_forwarded = request.META.get('HTTP_X_FORWARDED_FOR')
    if x_forwarded:
        return x_forwarded.split(',')[0].strip()
    return request.META.get('REMOTE_ADDR', '0.0.0.0')


def get_client_identifier(request):
    ip = get_client_ip(request)
    ua = request.META.get('HTTP_USER_AGENT', 'unknown')[:200]
    return f"{ip}|{ua}"


def get_session_key(request):
    if not request.session.session_key:
        request.session.create()
    return request.session.session_key


# ---------------------------------------------------------------------------
# Форматирование JSON → читаемый HTML (для страниц сайта)
# ---------------------------------------------------------------------------

def _convert_json_to_html(data):
    """Превращает JSON в читаемый HTML. Никакого сырого JSON."""
    
    if isinstance(data, list):
        if data and isinstance(data[0], dict) and 'name' in data[0]:
            html = ''
            for i, stage in enumerate(data, 1):
                name = stage.get('name', f'Этап {i}')
                duration = stage.get('duration_minutes', '')
                stype = stage.get('type', '')
                stype_ru = TYPE_TRANSLATIONS.get(stype, stype)
                html += f'<h3>{i}. {name} ({duration} мин, {stype_ru})</h3>'
                
                if stage.get('description'):
                    html += f'<p>{stage["description"]}</p>'
                
                if stage.get('teacher_actions'):
                    html += '<p><strong>Действия учителя:</strong></p><ul>'
                    for a in stage['teacher_actions']:
                        if isinstance(a, str):
                            html += f'<li>{a}</li>'
                        elif isinstance(a, dict):
                            html += f'<li>{a.get("action", str(a))}</li>'
                    html += '</ul>'
                
                if stage.get('student_actions'):
                    html += '<p><strong>Действия учеников:</strong></p><ul>'
                    for a in stage['student_actions']:
                        if isinstance(a, str):
                            html += f'<li>{a}</li>'
                        elif isinstance(a, dict):
                            html += f'<li>{a.get("action", str(a))}</li>'
                    html += '</ul>'
                
                if stage.get('questions'):
                    html += '<p><strong>Вопросы и предполагаемые ответы:</strong></p><ul>'
                    for q in stage['questions']:
                        if isinstance(q, str):
                            html += f'<li>{q}</li>'
                        elif isinstance(q, dict):
                            question_text = q.get('question', q.get('text', str(q)))
                            html += f'<li><strong>Вопрос:</strong> {question_text}'
                            if q.get('expected_answers'):
                                html += '<ul>'
                                for ans in q['expected_answers']:
                                    if isinstance(ans, str):
                                        html += f'<li><em>Возможный ответ:</em> {ans}</li>'
                                    elif isinstance(ans, dict):
                                        html += f'<li><em>{ans.get("student", "Ученик")}:</em> {ans.get("answer", str(ans))}</li>'
                                html += '</ul>'
                            if q.get('teacher_response'):
                                html += f'<p><em>Реакция учителя:</em> {q["teacher_response"]}</p>'
                            html += '</li>'
                    html += '</ul>'
                
                if stage.get('materials'):
                    html += '<p><strong>Материалы:</strong></p><ul>'
                    for m in stage['materials']:
                        if isinstance(m, str):
                            html += f'<li>{m}</li>'
                        elif isinstance(m, dict):
                            html += f'<li>{m.get("name", m.get("item", str(m)))}'
                            if m.get('quantity'):
                                html += f' — {m["quantity"]}'
                            html += '</li>'
                    html += '</ul>'
                
                if stage.get('mechanics'):
                    html += f'<p><strong>Механика:</strong> {stage["mechanics"]}</p>'
                
                if stage.get('expected_result'):
                    html += f'<p><strong>Ожидаемый результат:</strong> {stage["expected_result"]}</p>'
                
                html += '<hr>'
            return html
        
        elif data and isinstance(data[0], dict):
            html = ''
            for i, item in enumerate(data, 1):
                html += f'<h4>{i}.</h4>'
                for key, value in item.items():
                    key_ru = translate_key(key)
                    if isinstance(value, list):
                        html += f'<p><strong>{key_ru}:</strong></p><ul>'
                        for v in value:
                            html += f'<li>{v if isinstance(v, str) else json.dumps(v, ensure_ascii=False)}</li>'
                        html += '</ul>'
                    elif isinstance(value, dict):
                        html += f'<p><strong>{key_ru}:</strong></p>'
                        for k, v in value.items():
                            html += f'<p>{translate_key(k)}: {v}</p>'
                    else:
                        html += f'<p><strong>{key_ru}:</strong> {value}</p>'
                if i < len(data):
                    html += '<hr>'
            return html
        
        else:
            html = '<ol>'
            for item in data:
                if isinstance(item, str):
                    html += f'<li>{item}</li>'
                elif isinstance(item, dict):
                    html += '<li>'
                    for k, v in item.items():
                        html += f'<strong>{translate_key(k)}:</strong> {v}<br>'
                    html += '</li>'
                else:
                    html += f'<li>{item}</li>'
            html += '</ol>'
            return html
    
    elif isinstance(data, dict):
        html = ''
        for key, value in data.items():
            key_ru = translate_key(key)
            
            if isinstance(value, list):
                html += f'<p><strong>{key_ru}:</strong></p><ul>'
                for item in value:
                    if isinstance(item, str):
                        html += f'<li>{item}</li>'
                    elif isinstance(item, dict):
                        html += '<li>'
                        for k, v in item.items():
                            html += f'<strong>{translate_key(k)}:</strong> {v}<br>'
                        html += '</li>'
                html += '</ul>'
            
            elif isinstance(value, dict):
                html += f'<p><strong>{key_ru}:</strong></p>'
                for k, v in value.items():
                    k_ru = translate_key(k)
                    if isinstance(v, list):
                        html += f'<p><em>{k_ru}:</em></p><ul>'
                        for item in v:
                            html += f'<li>{item}</li>'
                        html += '</ul>'
                    else:
                        html += f'<p><em>{k_ru}:</em> {v}</p>'
            
            else:
                html += f'<p><strong>{key_ru}:</strong> {value}</p>'
        
        return html
    
    return str(data)


# ---------------------------------------------------------------------------
# Форматирование JSON → читаемый текст (для экспорта PDF/DOCX)
# ---------------------------------------------------------------------------

def parse_json_field(value):
    if not value or not str(value).strip():
        return ''
    value = str(value).strip()
    try:
        data = json.loads(value)
        return _format_json_data(data)
    except (json.JSONDecodeError, TypeError):
        return value


def _format_json_data(data):
    """Рекурсивно форматирует JSON в читаемый текст. Все ключи — на русском."""
    if isinstance(data, list):
        lines = []
        for i, item in enumerate(data, 1):
            if isinstance(item, dict):
                if 'name' in item:
                    name = item.get('name', f'Этап {i}')
                    duration = item.get('duration_minutes', '')
                    stype = item.get('type', '')
                    stype_ru = TYPE_TRANSLATIONS.get(stype, stype)
                    lines.append(f'{i}. {name} ({duration} мин, {stype_ru})')
                    for key, value in item.items():
                        if key in ('name', 'duration_minutes', 'type'):
                            continue
                        key_ru = translate_key(key)
                        if isinstance(value, list):
                            lines.append(f'   {key_ru}:')
                            for v in value:
                                if isinstance(v, str):
                                    lines.append(f'      - {v}')
                                elif isinstance(v, dict):
                                    if 'question' in v:
                                        lines.append(f'      Вопрос: {v.get("question", "")}')
                                        if v.get('expected_answers'):
                                            lines.append(f'         Ожидаемые ответы:')
                                            for ans in v['expected_answers']:
                                                if isinstance(ans, str):
                                                    lines.append(f'            - {ans}')
                                                elif isinstance(ans, dict):
                                                    lines.append(f'            - {ans.get("student", "")}: {ans.get("answer", "")}')
                                    else:
                                        lines.append(f'      - {json.dumps(v, ensure_ascii=False)}')
                        elif isinstance(value, str) and value.strip():
                            lines.append(f'   {key_ru}: {value}')
                    lines.append('')
                else:
                    for key, value in item.items():
                        key_ru = translate_key(key)
                        if isinstance(value, list):
                            lines.append(f'{key_ru}:')
                            for v in value:
                                lines.append(f'   - {v}')
                        elif isinstance(value, dict):
                            lines.append(f'{key_ru}:')
                            for k, v in value.items():
                                lines.append(f'   {translate_key(k)}: {v}')
                        else:
                            lines.append(f'{key_ru}: {value}')
            elif isinstance(item, str):
                lines.append(f'{i}. {item}')
            else:
                lines.append(f'{i}. {item}')
        return '\n'.join(lines)
    elif isinstance(data, dict):
        lines = []
        for key, value in data.items():
            key_ru = translate_key(key)
            if isinstance(value, list):
                lines.append(f'{key_ru}:')
                for item in value:
                    if isinstance(item, str):
                        lines.append(f'   - {item}')
                    elif isinstance(item, dict):
                        if 'question' in item:
                            lines.append(f'   Вопрос: {item.get("question", "")}')
                            if item.get('expected_answers'):
                                lines.append(f'      Ожидаемые ответы:')
                                for ans in item['expected_answers']:
                                    if isinstance(ans, str):
                                        lines.append(f'         - {ans}')
                                    elif isinstance(ans, dict):
                                        lines.append(f'         - {ans.get("student", "")}: {ans.get("answer", "")}')
                        else:
                            lines.append(f'   - {json.dumps(item, ensure_ascii=False)}')
                    else:
                        lines.append(f'   - {item}')
            elif isinstance(value, dict):
                lines.append(f'{key_ru}:')
                for k, v in value.items():
                    k_ru = translate_key(k)
                    if isinstance(v, list):
                        lines.append(f'   {k_ru}:')
                        for item in v:
                            lines.append(f'      - {item}')
                    else:
                        lines.append(f'   {k_ru}: {v}')
            elif isinstance(value, str):
                if value.strip():
                    lines.append(f'{key_ru}: {value}')
            else:
                lines.append(f'{key_ru}: {value}')
        return '\n'.join(lines)
    return str(data)


# ---------------------------------------------------------------------------
# Список сценариев
# ---------------------------------------------------------------------------

def scenario_list(request):
    session_key = get_session_key(request)
    scenarios = Scenario.objects.filter(session_key=session_key)
    return render(request, 'scenarios/list.html', {'scenarios': scenarios})


# ---------------------------------------------------------------------------
# Создание сценария
# ---------------------------------------------------------------------------

def scenario_create(request):
    if request.method == 'POST':
        form = ScenarioCreateForm(request.POST)
        if form.is_valid():
            scenario = form.save(commit=False)
            scenario.session_key = get_session_key(request)

            date_str = request.GET.get('date', '')
            if date_str:
                scenario.scheduled_date = datetime.strptime(date_str, '%Y-%m-%d').date()

            scenario.save()
            request.session['last_scenario_params'] = form.build_params_for_prompt()
            request.session['last_scenario_id'] = scenario.pk
            return redirect('scenario-generate-auto', pk=scenario.pk)
        else:
            messages.error(request, 'Пожалуйста, исправьте ошибки в форме.')
    else:
        form = ScenarioCreateForm()

    today = date.today()
    client_id = get_client_identifier(request)
    llm_requests_today = LLMRequestLog.objects.filter(
        client_identifier=client_id,
        created_at__date=today
    ).count()
    max_requests = 20

    return render(request, 'scenarios/create.html', {
        'form': form,
        'requests_used': llm_requests_today,
        'requests_max': max_requests,
        'requests_remaining': max_requests - llm_requests_today,
        'limit_exceeded': llm_requests_today >= max_requests,
    })


# ---------------------------------------------------------------------------
# Страницы генерации
# ---------------------------------------------------------------------------

def scenario_generate(request, pk):
    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk, session_key=session_key)
    params = request.session.get('last_scenario_params', {})
    return render(request, 'scenarios/generate.html', {
        'scenario': scenario,
        'params': params,
    })


def scenario_generate_auto(request, pk):
    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk, session_key=session_key)
    params = request.session.get('last_scenario_params', {})
    return render(request, 'scenarios/generate.html', {
        'scenario': scenario,
        'params': params,
        'auto_start': True,
    })


# ---------------------------------------------------------------------------
# API генерации
# ---------------------------------------------------------------------------

def scenario_generate_api(request, pk):
    if request.method != 'POST':
        return JsonResponse({'error': 'Method not allowed'}, status=405)

    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk, session_key=session_key)
    params = request.session.get('last_scenario_params', {})

    if not params:
        return JsonResponse({'error': 'Параметры генерации не найдены.'}, status=400)

    client_id = get_client_identifier(request)
    today = date.today()
    llm_requests_today = LLMRequestLog.objects.filter(
        client_identifier=client_id,
        created_at__date=today
    ).count()

    max_requests = 20
    if llm_requests_today >= max_requests:
        return JsonResponse({
            'error': f'Дневной лимит запросов исчерпан ({max_requests}/{max_requests}). Попробуйте завтра.',
            'limit_exceeded': True,
            'requests_used': llm_requests_today,
            'requests_max': max_requests,
        }, status=429)

    try:
        service = ScenarioGeneratorService()
        result = service.generate_scenario(params, session_key=session_key)

        scenario.title = result.get('title', '')
        scenario.legend = result.get('legend', '')
        scenario.goals = json.dumps(result.get('goals', {}), ensure_ascii=False, indent=2)
        scenario.stages = json.dumps(result.get('stages', []), ensure_ascii=False, indent=2)
        scenario.materials = json.dumps(result.get('materials_full', {}), ensure_ascii=False, indent=2)
        scenario.adaptation = json.dumps(result.get('adaptation', {}), ensure_ascii=False, indent=2)
        scenario.reflection = json.dumps(result.get('reflection', {}), ensure_ascii=False, indent=2)
        scenario.risks = json.dumps(result.get('risks', {}), ensure_ascii=False, indent=2)
        scenario.assessment = json.dumps(result.get('assessment', {}), ensure_ascii=False, indent=2)
        scenario.teacher_tips = json.dumps(result.get('teacher_tips', []), ensure_ascii=False, indent=2)
        scenario.host_script = result.get('host_script', '')
        scenario.checklist = json.dumps(result.get('checklist', {}), ensure_ascii=False, indent=2)
        scenario.homework = json.dumps(result.get('homework', {}), ensure_ascii=False, indent=2)
        scenario.save()

        LLMRequestLog.objects.create(
            session_key=session_key,
            client_identifier=client_id,
            request_type='generate_scenario',
            tokens_used=0,
        )

        version_service = ScenarioVersionService(scenario)
        version_service.create_version(
            change_description='Автоматическая генерация сценария',
            created_by=session_key,
        )

        # Сохраняем эмбеддинг для RAG
        try:
            from .embedding_service import EmbeddingService
            emb_service = EmbeddingService()
            emb_service.embed_scenario(scenario)
        except Exception as e:
            logger.warning(f"⚠️ Не удалось создать эмбеддинг: {e}")

        return JsonResponse({
            'status': 'success',
            'result': result,
            'requests_used': llm_requests_today + 1,
            'requests_max': max_requests,
        })

    except ValueError as e:
        return JsonResponse({'error': str(e)}, status=500)
    except Exception as e:
        logger.error(f'Ошибка генерации: {e}', exc_info=True)
        return JsonResponse({'error': f'Ошибка генерации: {str(e)}'}, status=500)


# ---------------------------------------------------------------------------
# Перегенерация блока
# ---------------------------------------------------------------------------

def scenario_regenerate_block_api(request, pk, block_index):
    if request.method != 'POST':
        return JsonResponse({'error': 'Method not allowed'}, status=405)

    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk, session_key=session_key)

    try:
        stages = json.loads(scenario.stages)
        if block_index < 0 or block_index >= len(stages):
            return JsonResponse({'error': 'Неверный индекс блока'}, status=400)

        current_block = stages[block_index]
        block_type = current_block.get('type', 'main')

        full_context = {
            'title': scenario.title,
            'format_type': scenario.get_format_type_display(),
            'grade': scenario.grade,
            'theme': scenario.theme,
            'duration': scenario.duration,
        }

        service = ScenarioGeneratorService()
        new_block = service.regenerate_block(block_type, current_block, full_context)

        stages[block_index] = new_block
        scenario.stages = json.dumps(stages, ensure_ascii=False, indent=2)
        scenario.save()

        version_service = ScenarioVersionService(scenario)
        version_service.create_version(
            change_description=f'Перегенерирован этап: {new_block.get("name", block_index + 1)}',
            created_by=session_key,
        )

        return JsonResponse({'status': 'success', 'block': new_block})

    except Exception as e:
        logger.error(f'Ошибка перегенерации: {e}', exc_info=True)
        return JsonResponse({'error': f'Ошибка перегенерации: {str(e)}'}, status=500)


# ---------------------------------------------------------------------------
# Детальный просмотр
# ---------------------------------------------------------------------------

def scenario_detail(request, pk):
    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk)

    url_key = request.GET.get('key', '')
    has_access = (scenario.session_key == session_key) or (scenario.share_key and scenario.share_key == url_key)

    if not has_access:
        return render(request, 'scenarios/access_denied.html', status=403)

    def to_html(value):
        if not value or not value.strip():
            return ''
        value = value.strip()
        if value.startswith('<'):
            return value
        try:
            data = json.loads(value)
            return _convert_json_to_html(data)
        except (json.JSONDecodeError, TypeError):
            return value.replace('\n', '<br>')

    return render(request, 'scenarios/detail.html', {
        'scenario': scenario,
        'legend': to_html(scenario.legend),
        'goals': to_html(scenario.goals),
        'stages': to_html(scenario.stages),
        'materials': to_html(scenario.materials),
        'adaptation': to_html(scenario.adaptation),
        'reflection': to_html(scenario.reflection),
        'risks': to_html(scenario.risks),
        'assessment': to_html(scenario.assessment),
        'teacher_tips': to_html(scenario.teacher_tips),
        'host_script': to_html(scenario.host_script),
        'checklist': to_html(scenario.checklist),
        'homework': to_html(scenario.homework),
    })


# ---------------------------------------------------------------------------
# Редактирование
# ---------------------------------------------------------------------------

def scenario_edit(request, pk):
    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk, session_key=session_key)

    if request.method == 'POST':
        scenario.title = request.POST.get('title', '').strip()
        scenario.legend = request.POST.get('legend', '').strip()
        scenario.goals = request.POST.get('goals', '').strip()
        scenario.stages = request.POST.get('stages', '').strip()
        scenario.materials = request.POST.get('materials', '').strip()
        scenario.adaptation = request.POST.get('adaptation', '').strip()
        scenario.reflection = request.POST.get('reflection', '').strip()
        scenario.risks = request.POST.get('risks', '').strip()
        scenario.assessment = request.POST.get('assessment', '').strip()
        scenario.teacher_tips = request.POST.get('teacher_tips', '').strip()
        scenario.host_script = request.POST.get('host_script', '').strip()
        scenario.checklist = request.POST.get('checklist', '').strip()
        scenario.homework = request.POST.get('homework', '').strip()
        scenario.save()

        version_service = ScenarioVersionService(scenario)
        change_desc = request.POST.get('change_description', '').strip()
        version_service.create_version(
            change_description=change_desc or 'Ручное редактирование',
            created_by=session_key,
        )

        messages.success(request, 'Сценарий обновлён!')
        return redirect('scenario-edit', pk=scenario.pk)

    def to_html(value):
        if not value or not value.strip():
            return ''
        value = value.strip()
        if value.startswith('<'):
            return value
        try:
            data = json.loads(value)
            return _convert_json_to_html(data)
        except (json.JSONDecodeError, TypeError):
            return value

    return render(request, 'scenarios/edit.html', {
        'scenario': scenario,
        'title': scenario.title or '',
        'legend': to_html(scenario.legend) or scenario.legend or '',
        'goals': to_html(scenario.goals) or scenario.goals or '',
        'stages': to_html(scenario.stages) or scenario.stages or '',
        'materials': to_html(scenario.materials) or scenario.materials or '',
        'adaptation': to_html(scenario.adaptation) or scenario.adaptation or '',
        'reflection': to_html(scenario.reflection) or scenario.reflection or '',
        'risks': to_html(scenario.risks) or scenario.risks or '',
        'assessment': to_html(scenario.assessment) or scenario.assessment or '',
        'teacher_tips': to_html(scenario.teacher_tips) or scenario.teacher_tips or '',
        'host_script': scenario.host_script or '',
        'checklist': to_html(scenario.checklist) or scenario.checklist or '',
        'homework': to_html(scenario.homework) or scenario.homework or '',
    })


# ---------------------------------------------------------------------------
# Удаление
# ---------------------------------------------------------------------------

def scenario_delete(request, pk):
    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk, session_key=session_key)

    if request.method == 'POST':
        scenario.delete()
        messages.success(request, 'Сценарий удалён')
        return redirect('scenario-list')

    return render(request, 'scenarios/delete.html', {'scenario': scenario})


# ---------------------------------------------------------------------------
# Шеринг
# ---------------------------------------------------------------------------

def scenario_share(request, pk):
    import uuid
    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk, session_key=session_key)

    if not scenario.share_key:
        scenario.share_key = uuid.uuid4().hex[:16]
        scenario.save()

    share_url = request.build_absolute_uri(f'/scenarios/{pk}/?key={scenario.share_key}')
    return JsonResponse({'status': 'success', 'url': share_url})


def scenario_unshare(request, pk):
    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk, session_key=session_key)
    scenario.share_key = None
    scenario.save()
    return JsonResponse({'status': 'success'})


# ---------------------------------------------------------------------------
# Дата проведения
# ---------------------------------------------------------------------------

def scenario_set_date(request, pk):
    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk, session_key=session_key)

    if request.method == 'POST':
        date_str = request.POST.get('date', '')
        if date_str:
            scenario.scheduled_date = datetime.strptime(date_str, '%Y-%m-%d').date()
            scenario.save()
            return JsonResponse({'status': 'success', 'message': 'Дата сохранена'})
        return JsonResponse({'status': 'error', 'message': 'Дата не указана'}, status=400)

    return JsonResponse({'status': 'error'}, status=405)


def scenario_remove_date(request, pk):
    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk, session_key=session_key)

    if request.method == 'POST':
        scenario.scheduled_date = None
        scenario.save()
        return JsonResponse({'status': 'success'})

    return JsonResponse({'status': 'error'}, status=405)


# ---------------------------------------------------------------------------
# QR-код
# ---------------------------------------------------------------------------

def scenario_qrcode(request, pk):
    import qrcode

    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk, session_key=session_key)

    if not scenario.share_key:
        import uuid
        scenario.share_key = uuid.uuid4().hex[:16]
        scenario.save()

    share_url = request.build_absolute_uri(f'/scenarios/{pk}/?key={scenario.share_key}')

    qr = qrcode.QRCode(box_size=10, border=4)
    qr.add_data(share_url)
    qr.make(fit=True)
    img = qr.make_image(fill_color='#2563eb', back_color='white')

    buf = io.BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)

    return HttpResponse(buf.read(), content_type='image/png')


# ---------------------------------------------------------------------------
# Календарь: загрузка плана
# ---------------------------------------------------------------------------

def calendar_upload_plan(request):
    session_key = get_session_key(request)

    if request.method != 'POST':
        return JsonResponse({'status': 'error', 'error': 'Method not allowed'}, status=405)

    text = request.POST.get('plan_text', '').strip()
    file_type = 'txt'

    if not text and 'plan_file' in request.FILES:
        file = request.FILES['plan_file']
        file_name = file.name.lower()

        try:
            if file_name.endswith('.txt'):
                text = file.read().decode('utf-8')
                file_type = 'txt'
            elif file_name.endswith('.csv'):
                text = file.read().decode('utf-8')
                file_type = 'csv'
            elif file_name.endswith('.docx'):
                text = _parse_docx(file)
                file_type = 'docx'
            elif file_name.endswith('.pdf'):
                text = _parse_pdf(file)
                file_type = 'pdf'
            elif file_name.endswith('.pptx'):
                text = _parse_pptx(file)
                file_type = 'pptx'
            else:
                return JsonResponse({
                    'status': 'error',
                    'error': 'Формат файла не поддерживается. Поддерживаются: TXT, CSV, DOCX, PDF, PPTX',
                }, status=400)
        except UnicodeDecodeError:
            file.seek(0)
            try:
                text = file.read().decode('cp1251')
                file_type = 'txt'
            except Exception as e:
                return JsonResponse({'status': 'error', 'error': f'Ошибка кодировки файла: {e}'}, status=400)
        except Exception as e:
            logger.error(f'Ошибка чтения файла: {e}')
            return JsonResponse({'status': 'error', 'error': f'Ошибка чтения файла: {e}'}, status=400)

    if not text:
        return JsonResponse({'status': 'error', 'error': 'Введите текст плана или загрузите файл'}, status=400)

    if len(text) > 100000:
        return JsonResponse({'status': 'error', 'error': 'Текст слишком большой. Максимум 100 000 символов.'}, status=400)

    try:
        from .services import DeepSeekClient

        client = DeepSeekClient()

        system_prompt = """Ты — парсер плана воспитательной работы. Извлеки из текста все даты и темы внеурочных занятий.

## ПРАВИЛА
1. Ищи даты в любом формате: "7 мая", "07.05", "2026-05-07", "12 числа", "в конце месяца"
2. Если год не указан — используй текущий 2026 год
3. Если нет конкретного дня месяца — используй 1 число
4. Класс ищи в форматах: "7 класс", "для 7 класса", "7 кл", "7А"
5. Если тема не указана — напиши "Тема не указана"
6. Если дата не указана — пропусти эту строку

## ФОРМАТ ОТВЕТА
Верни ТОЛЬКО JSON-массив:
[{"date": "YYYY-MM-DD", "theme": "Тема", "grade": число или null}]"""

        user_prompt = f"""Извлеки даты и темы из этого текста плана воспитательной работы:

{text}"""

        result = client.generate(
            system_prompt=system_prompt,
            user_input=user_prompt,
            temperature=0.1,
            max_tokens=2000,
            max_retries=3,
        )

        if isinstance(result, dict):
            themes = result.get('themes', result.get('data', result.get('result', [])))
            if isinstance(themes, dict):
                themes = [themes]
        elif isinstance(result, list):
            themes = result
        else:
            themes = []

        if not themes:
            return JsonResponse({
                'status': 'success', 'count': 0, 'skipped': 0,
                'total_found': 0, 'themes': [],
                'message': 'В тексте не найдено дат и тем',
                'file_type': file_type,
            })

        total_found = len(themes)
        added_count = 0
        skipped_count = 0
        added_themes = []
        skipped_themes = []

        for item in themes:
            if not isinstance(item, dict):
                skipped_count += 1
                continue

            date_str = item.get('date', '')
            theme_text = item.get('theme', 'Тема не указана')
            grade = item.get('grade', None)

            if not date_str or not theme_text:
                skipped_count += 1
                continue

            try:
                date_obj = datetime.strptime(date_str, '%Y-%m-%d').date()
            except ValueError:
                try:
                    date_obj = datetime.strptime(date_str, '%d.%m.%Y').date()
                except ValueError:
                    skipped_count += 1
                    continue

            duplicate_exists = ScheduledTheme.objects.filter(
                session_key=session_key,
                date=date_obj,
                theme__iexact=theme_text.strip(),
            ).exists()

            if duplicate_exists:
                skipped_count += 1
                skipped_themes.append({'date': date_str, 'theme': theme_text, 'reason': 'дубликат'})
                continue

            ScheduledTheme.objects.create(
                session_key=session_key,
                theme=theme_text.strip(),
                grade=grade if grade and str(grade).isdigit() else None,
                date=date_obj,
            )
            added_count += 1
            added_themes.append({'date': date_str, 'theme': theme_text, 'grade': grade})

        return JsonResponse({
            'status': 'success',
            'count': added_count,
            'skipped': skipped_count,
            'total_found': total_found,
            'themes': added_themes,
            'file_type': file_type,
            'message': f'Добавлено: {added_count}, пропущено: {skipped_count} (из них дубликатов: {len(skipped_themes)})',
        })

    except Exception as e:
        logger.error(f"Ошибка загрузки плана: {e}", exc_info=True)
        return JsonResponse({'status': 'error', 'error': str(e)}, status=500)


# ---------------------------------------------------------------------------
# Календарь: удаление темы
# ---------------------------------------------------------------------------

def calendar_delete_theme(request):
    if request.method != 'POST':
        return JsonResponse({'status': 'error', 'error': 'Method not allowed'}, status=405)

    session_key = get_session_key(request)
    theme_id = request.POST.get('theme_id', '').strip()

    if not theme_id:
        return JsonResponse({'status': 'error', 'error': 'Не указан идентификатор темы'}, status=400)

    try:
        theme_id_int = int(theme_id)
    except (ValueError, TypeError):
        return JsonResponse({'status': 'error', 'error': 'Некорректный идентификатор темы'}, status=400)

    try:
        theme = ScheduledTheme.objects.get(id=theme_id_int, session_key=session_key)
    except ScheduledTheme.DoesNotExist:
        return JsonResponse({'status': 'error', 'error': 'Тема не найдена или недоступна для удаления'}, status=404)

    theme_text = theme.theme
    theme_date = theme.date
    theme.delete()

    logger.info(f'Тема удалена из календаря: "{theme_text}" на {theme_date}. Сессия: {session_key[:10]}...')

    return JsonResponse({'status': 'success', 'message': f'Тема «{theme_text}» удалена из календаря'})


# ---------------------------------------------------------------------------
# Версионирование
# ---------------------------------------------------------------------------

def scenario_versions(request, pk):
    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk, session_key=session_key)

    version_service = ScenarioVersionService(scenario)
    versions = version_service.get_versions()

    return render(request, 'scenarios/versions.html', {
        'scenario': scenario,
        'versions': versions,
    })


def scenario_version_detail(request, pk, version_number):
    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk, session_key=session_key)

    version_service = ScenarioVersionService(scenario)
    version = version_service.get_version(version_number)

    if not version:
        messages.error(request, f'Версия {version_number} не найдена')
        return redirect('scenario-versions', pk=pk)

    return render(request, 'scenarios/version_detail.html', {
        'scenario': scenario,
        'version': version,
    })


def scenario_restore_version(request, pk, version_number):
    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk, session_key=session_key)

    if request.method != 'POST':
        return JsonResponse({'status': 'error', 'error': 'Method not allowed'}, status=405)

    try:
        version_service = ScenarioVersionService(scenario)
        version_service.restore_version(version_number)
        messages.success(request, f'Сценарий восстановлен до версии {version_number}')
        return JsonResponse({'status': 'success'})
    except ValueError as e:
        return JsonResponse({'status': 'error', 'error': str(e)}, status=404)
    except Exception as e:
        logger.error(f'Ошибка восстановления версии: {e}', exc_info=True)
        return JsonResponse({'status': 'error', 'error': str(e)}, status=500)


def scenario_compare_versions(request, pk):
    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk, session_key=session_key)

    v1 = request.GET.get('v1')
    v2 = request.GET.get('v2')

    if not v1 or not v2:
        return JsonResponse({'status': 'error', 'error': 'Укажите v1 и v2'}, status=400)

    try:
        version_service = ScenarioVersionService(scenario)
        comparison = version_service.compare_versions(int(v1), int(v2))
        return JsonResponse({'status': 'success', 'comparison': comparison})
    except ValueError as e:
        return JsonResponse({'status': 'error', 'error': str(e)}, status=404)
    except Exception as e:
        logger.error(f'Ошибка сравнения версий: {e}', exc_info=True)
        return JsonResponse({'status': 'error', 'error': str(e)}, status=500)


def scenario_delete_version(request, pk, version_number):
    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk, session_key=session_key)

    if request.method != 'POST':
        return JsonResponse({'status': 'error', 'error': 'Method not allowed'}, status=405)

    version_service = ScenarioVersionService(scenario)
    result = version_service.delete_version(version_number)

    if result['status'] == 'error':
        return JsonResponse({'status': 'error', 'error': result['message']}, status=404)

    if result['status'] == 'scenario_deleted':
        messages.success(request, result['message'])
        return JsonResponse({
            'status': 'scenario_deleted',
            'message': result['message'],
            'redirect': '/scenarios/',
        })

    messages.success(request, result['message'])
    return JsonResponse({'status': 'success', 'message': result['message']})


# ---------------------------------------------------------------------------
# Статистика
# ---------------------------------------------------------------------------

def scenario_statistics(request):
    session_key = get_session_key(request)
    client_id = get_client_identifier(request)
    scenarios = Scenario.objects.filter(session_key=session_key)

    total_count = scenarios.count()

    direction_counts = {}
    for direction_code, direction_label in Scenario.DIRECTION_CHOICES:
        count = scenarios.filter(direction=direction_code).count()
        direction_counts[direction_label] = count

    format_counts = {}
    for format_code, format_label in Scenario.FORMAT_CHOICES:
        count = scenarios.filter(format_type=format_code).count()
        format_counts[format_label] = count

    grade_counts = {}
    for grade in range(1, 12):
        count = scenarios.filter(grade=grade).count()
        grade_counts[grade] = count

    month_counts = {str(i): 0 for i in range(1, 13)}
    month_names = {
        '1': 'Янв', '2': 'Фев', '3': 'Мар', '4': 'Апр',
        '5': 'Май', '6': 'Июн', '7': 'Июл', '8': 'Авг',
        '9': 'Сен', '10': 'Окт', '11': 'Ноя', '12': 'Дек',
    }
    for scenario in scenarios:
        month_key = str(scenario.created_at.month)
        month_counts[month_key] = month_counts.get(month_key, 0) + 1

    this_month_count = scenarios.filter(created_at__month=date.today().month).count()

    days_with_events = Scenario.objects.filter(
        session_key=session_key,
        scheduled_date__isnull=False,
    ).values('scheduled_date').distinct().count()

    top_themes = scenarios.values('theme').annotate(
        count=Count('id')
    ).order_by('-count')[:5]

    llm_requests = LLMRequestLog.objects.filter(client_identifier=client_id).count()
    llm_requests_today = LLMRequestLog.objects.filter(
        client_identifier=client_id,
        created_at__date=date.today(),
    ).count()

    time_saved = total_count * 2.5

    today = date.today()
    days_in_month = 30
    scheduled_days = Scenario.objects.filter(
        session_key=session_key,
        scheduled_date__gte=today,
        scheduled_date__lte=today + timedelta(days=30),
    ).values('scheduled_date').distinct().count()
    calendar_fill = round((scheduled_days / days_in_month) * 100) if days_in_month > 0 else 0

    max_requests = 20
    requests_remaining = max(0, max_requests - llm_requests_today)

    context = {
        'total_count': total_count,
        'this_month_count': this_month_count,
        'days_with_events': days_with_events,
        'time_saved': int(time_saved),
        'direction_counts': json.dumps(direction_counts),
        'format_counts': json.dumps(format_counts),
        'grade_counts': json.dumps(grade_counts),
        'month_counts': json.dumps(month_counts),
        'month_names': json.dumps(month_names),
        'top_themes': top_themes,
        'llm_requests': llm_requests,
        'llm_requests_today': llm_requests_today,
        'calendar_fill': calendar_fill,
        'scheduled_days': scheduled_days,
        'requests_remaining': requests_remaining,
        'requests_max': max_requests,
    }

    return render(request, 'scenarios/statistics.html', context)


# ---------------------------------------------------------------------------
# Экспорт в PDF
# ---------------------------------------------------------------------------

def scenario_export_pdf(request, pk):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
    from reportlab.lib import colors
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from html.parser import HTMLParser

    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk, session_key=session_key)

    class HTMLToText(HTMLParser):
        def __init__(self):
            super().__init__()
            self.result = []
            self.skip = False
        def handle_starttag(self, tag, attrs):
            if tag in ('h1','h2','h3'):
                self.result.append('\n')
            elif tag == 'li':
                self.result.append('\n• ')
            elif tag in ('script','style'):
                self.skip = True
        def handle_endtag(self, tag):
            if tag in ('h1','h2','h3','p','ul','ol','div'):
                self.result.append('\n')
            elif tag in ('script','style'):
                self.skip = False
        def handle_data(self, data):
            if not self.skip:
                self.result.append(data.strip())

    def strip_html(text):
        if not text:
            return ''
        parser = HTMLToText()
        parser.feed(text)
        return ' '.join(parser.result).strip()

    font_path = os.path.join(settings.BASE_DIR, 'static', 'fonts', 'DejaVuSans.ttf')
    font_bold_path = os.path.join(settings.BASE_DIR, 'static', 'fonts', 'DejaVuSans-Bold.ttf')

    if os.path.exists(font_path):
        pdfmetrics.registerFont(TTFont('DejaVu', font_path))
        fn = 'DejaVu'
    else:
        fn = 'Helvetica'

    if os.path.exists(font_bold_path):
        pdfmetrics.registerFont(TTFont('DejaVuBd', font_bold_path))
        fn_bold = 'DejaVuBd'
    else:
        fn_bold = fn

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=20*mm, rightMargin=15*mm,
        topMargin=20*mm, bottomMargin=20*mm
    )
    doc.title = scenario.title or scenario.theme

    text_color = colors.HexColor('#0f172a')
    muted_color = colors.HexColor('#64748b')

    title_style = ParagraphStyle(
        'Title', fontName=fn_bold, fontSize=20,
        spaceAfter=4, textColor=text_color, leading=26
    )
    h2_style = ParagraphStyle(
        'H2', fontName=fn_bold, fontSize=15,
        spaceBefore=16, spaceAfter=8, textColor=text_color, leading=20
    )
    normal = ParagraphStyle(
        'Normal', fontName=fn, fontSize=11,
        spaceAfter=4, leading=17, textColor=text_color
    )
    small = ParagraphStyle(
        'Small', fontName=fn, fontSize=9,
        textColor=muted_color, leading=13
    )

    story = []

    title_text = scenario.title or scenario.theme
    story.append(Paragraph(title_text, title_style))
    story.append(Paragraph(
        f"{scenario.get_direction_display()} | "
        f"{scenario.get_format_type_display()} | "
        f"{scenario.grade} класс | {scenario.duration} мин.",
        small
    ))
    story.append(Paragraph(f"<b>Тема:</b> {scenario.theme}", normal))
    story.append(Spacer(1, 8))

    sections = [
        ('Легенда', scenario.legend),
        ('Цели', scenario.goals),
        ('Этапы', scenario.stages),
        ('Материалы и ресурсы', scenario.materials),
        ('Адаптация', scenario.adaptation),
        ('Рефлексия', scenario.reflection),
        ('Риски', scenario.risks),
        ('Оценка эффективности', scenario.assessment),
        ('Советы учителю', scenario.teacher_tips),
        ('Сценарий для ведущего', scenario.host_script),
        ('Чек-лист подготовки', scenario.checklist),
        ('Домашнее задание', scenario.homework),
    ]

    for section_title, content in sections:
        if content and str(content).strip():
            story.append(Paragraph(section_title, h2_style))
            parsed = parse_json_field(content)
            for line in parsed.split('\n'):
                line = line.strip()
                if line:
                    safe_line = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    story.append(Paragraph(safe_line, normal))
            story.append(Spacer(1, 6))

    if scenario.share_key:
        import qrcode
        qr = qrcode.QRCode(box_size=4, border=2)
        qr.add_data(f'http://127.0.0.1:8000/scenarios/{pk}/?key={scenario.share_key}')
        qr.make(fit=True)
        qr_img = qr.make_image(fill_color='#2563eb', back_color='white')
        qr_buf = io.BytesIO()
        qr_img.save(qr_buf, format='PNG')
        qr_buf.seek(0)
        qr_image = Image(qr_buf, width=100, height=100)
        story.append(Spacer(1, 12))
        story.append(Paragraph("QR-код для доступа к сценарию:", h2_style))
        story.append(qr_image)

    doc.build(story)
    buf.seek(0)

    response = HttpResponse(buf.read(), content_type='application/pdf')
    filename = f"stsenariy_{scenario.pk}.pdf"
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    return response


# ---------------------------------------------------------------------------
# Экспорт в DOCX
# ---------------------------------------------------------------------------

def scenario_export_docx(request, pk):
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from html.parser import HTMLParser

    session_key = get_session_key(request)
    scenario = get_object_or_404(Scenario, pk=pk, session_key=session_key)

    class HTMLToText(HTMLParser):
        def __init__(self):
            super().__init__()
            self.result = []
            self.skip = False
        def handle_starttag(self, tag, attrs):
            if tag in ('h1','h2','h3'):
                self.result.append('\n')
            elif tag == 'li':
                self.result.append('\n• ')
            elif tag in ('script','style'):
                self.skip = True
        def handle_endtag(self, tag):
            if tag in ('h1','h2','h3','p','ul','ol','div'):
                self.result.append('\n')
            elif tag in ('script','style'):
                self.skip = False
        def handle_data(self, data):
            if not self.skip:
                self.result.append(data.strip())

    def strip_html(text):
        if not text:
            return ''
        parser = HTMLToText()
        parser.feed(text)
        return ' '.join(parser.result).strip()

    doc = Document()
    style = doc.styles['Normal']
    style.font.name = 'Arial'
    style.font.size = Pt(11)

    title = doc.add_heading(scenario.title or scenario.theme, level=1)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = meta.add_run(
        f"{scenario.get_direction_display()} | "
        f"{scenario.get_format_type_display()} | "
        f"{scenario.grade} класс | {scenario.duration} мин."
    )
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(100, 116, 139)

    p = doc.add_paragraph()
    run = p.add_run(f"Тема: {scenario.theme}")
    run.bold = True

    sections = [
        ('Легенда', scenario.legend),
        ('Цели', scenario.goals),
        ('Этапы', scenario.stages),
        ('Материалы и ресурсы', scenario.materials),
        ('Адаптация', scenario.adaptation),
        ('Рефлексия', scenario.reflection),
        ('Риски', scenario.risks),
        ('Оценка эффективности', scenario.assessment),
        ('Советы учителю', scenario.teacher_tips),
        ('Сценарий для ведущего', scenario.host_script),
        ('Чек-лист подготовки', scenario.checklist),
        ('Домашнее задание', scenario.homework),
    ]

    for section_title, content in sections:
        if content and str(content).strip():
            doc.add_heading(section_title, level=2)
            parsed = parse_json_field(content)
            for line in parsed.split('\n'):
                line = line.strip()
                if line:
                    doc.add_paragraph(line)

    if scenario.share_key:
        import qrcode
        qr = qrcode.QRCode(box_size=4, border=2)
        qr.add_data(f'http://127.0.0.1:8000/scenarios/{pk}/?key={scenario.share_key}')
        qr.make(fit=True)
        qr_img = qr.make_image(fill_color='#2563eb', back_color='white')
        qr_buf = io.BytesIO()
        qr_img.save(qr_buf, format='PNG')
        qr_buf.seek(0)
        doc.add_heading('QR-код для доступа к сценарию', level=2)
        doc.add_picture(qr_buf, width=Inches(1.5))

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    response = HttpResponse(
        buf.read(),
        content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
    )
    filename = f"stsenariy_{scenario.pk}.docx"
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    return response


# ---------------------------------------------------------------------------
# Парсеры файлов
# ---------------------------------------------------------------------------

def _parse_docx(file):
    from docx import Document
    doc = Document(file)
    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    return '\n'.join(paragraphs)


def _parse_pdf(file):
    from PyPDF2 import PdfReader
    reader = PdfReader(file)
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text and text.strip():
            pages.append(text.strip())
    return '\n'.join(pages)


def _parse_pptx(file):
    from pptx import Presentation
    prs = Presentation(file)
    slides = []
    for slide in prs.slides:
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        slide_texts.append(paragraph.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        slide_texts.append(row_text)
        if slide_texts:
            slides.append('\n'.join(slide_texts))
    return '\n\n'.join(slides)